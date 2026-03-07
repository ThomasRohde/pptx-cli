from __future__ import annotations

import json
import shutil
import struct
import zlib
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from typer.testing import CliRunner

from pptx_cli import __version__
from pptx_cli.cli import app

runner = CliRunner()


def _invoke_json(arguments: list[str]) -> dict[str, Any]:
    result = runner.invoke(app, arguments + ["--format", "json"])
    assert result.exit_code == 0, result.stdout
    return json.loads(result.stdout)


def _init_manifest(template_path: Path, output_dir: Path) -> dict[str, Any]:
    return _invoke_json(["init", str(template_path), "--out", str(output_dir)])


def _write_png(path: Path, width: int, height: int) -> None:
    def _chunk(tag: bytes, payload: bytes) -> bytes:
        checksum = zlib.crc32(tag + payload) & 0xFFFFFFFF
        return struct.pack(">I", len(payload)) + tag + payload + struct.pack(">I", checksum)

    row = bytes((23, 79, 129)) * width
    scanlines = b"".join(b"\x00" + row for _ in range(height))
    payload = b"".join(
        [
            b"\x89PNG\r\n\x1a\n",
            _chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)),
            _chunk(b"IDAT", zlib.compress(scanlines)),
            _chunk(b"IEND", b""),
        ]
    )
    path.write_bytes(payload)


def _placeholder_by_idx(slide: Any, placeholder_idx: int) -> Any:
    return next(
        shape for shape in slide.placeholders if shape.placeholder_format.idx == placeholder_idx
    )


def test_guide_returns_structured_envelope_with_catalog() -> None:
    payload = _invoke_json(["guide"])

    assert payload["ok"] is True
    assert payload["command"] == "guide.show"
    command_ids = {command["id"] for command in payload["result"]["commands"]}
    assert "layouts.list" in command_ids
    assert "deck.build" in command_ids
    assert payload["result"]["exit_codes"]["validation_error"] == 10
    assert "ERR_IO_NOT_FOUND" in payload["result"]["error_codes"]
    assert payload["result"]["content_objects"]["image"]["example"]["kind"] == "image"
    assert payload["result"]["content_objects"]["table"]["example"]["kind"] == "table"


def test_init_dry_run_returns_change_plan(repo_root: Path, template_path: Path) -> None:
    manifest_dir = repo_root / "tmp-manifest-dry-run"
    result = runner.invoke(
        app,
        [
            "init",
            str(template_path),
            "--out",
            str(manifest_dir),
            "--dry-run",
            "--format",
            "json",
        ],
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    assert payload["ok"] is True
    assert payload["result"]["dry_run"] is True
    assert any(
        change["target"].endswith("manifest.yaml")
        for change in payload["result"]["plan"]["changes"]
    )


def test_init_writes_manifest_package(template_path: Path, manifest_dir: Path) -> None:
    payload = _init_manifest(template_path, manifest_dir)

    assert payload["ok"] is True
    assert (manifest_dir / "manifest.yaml").exists()
    assert (manifest_dir / "manifest.schema.json").exists()
    assert (manifest_dir / "annotations.yaml").exists()
    assert (manifest_dir / "reports" / "init-report.json").exists()
    manifest_payload = yaml.safe_load((manifest_dir / "manifest.yaml").read_text(encoding="utf-8"))
    assert manifest_payload["template"]["stored_template_path"] == "assets/source-template.pptx"
    assert len(manifest_payload["layouts"]) >= 20
    picture_placeholder = next(
        placeholder
        for layout in manifest_payload["layouts"]
        for placeholder in layout["placeholders"]
        if placeholder["placeholder_type"] == "picture" and placeholder["logical_name"] == "picture"
    )
    assert picture_placeholder["supported_content_types"] == ["image"]


def test_init_fails_when_template_missing() -> None:
    result = runner.invoke(
        app,
        [
            "init",
            "missing-template.pptx",
            "--out",
            "./tmp-manifest",
            "--format",
            "json",
        ],
    )

    assert result.exit_code == 50
    payload = json.loads(result.stdout)
    assert payload["ok"] is False
    assert payload["errors"][0]["code"] == "ERR_IO_NOT_FOUND"


def test_version_flag_prints_installed_version() -> None:
    result = runner.invoke(app, ["--version"])

    assert result.exit_code == 0
    assert result.stdout.strip() == __version__


def test_inspection_commands_use_manifest_contract(
    template_path: Path,
    manifest_dir: Path,
) -> None:
    _init_manifest(template_path, manifest_dir)

    layouts_payload = _invoke_json(["layouts", "list", "--manifest", str(manifest_dir)])
    assert layouts_payload["result"]["count"] >= 20
    layout_ids = {layout["id"] for layout in layouts_payload["result"]["layouts"]}
    assert "title-only" in layout_ids

    layout_payload = _invoke_json(
        ["layouts", "show", "title-only", "--manifest", str(manifest_dir)]
    )
    assert layout_payload["result"]["id"] == "title-only"

    placeholders_payload = _invoke_json(
        ["placeholders", "list", "title-only", "--manifest", str(manifest_dir)]
    )
    placeholders = placeholders_payload["result"]["placeholders"]
    placeholder_names = {placeholder["logical_name"] for placeholder in placeholders}
    assert "title" in placeholder_names
    title_placeholder = next(
        placeholder for placeholder in placeholders if placeholder["logical_name"] == "title"
    )
    assert "Click to add slide title" in title_placeholder["guidance_text"]
    assert title_placeholder["text_defaults"]["max_lines"] == 2
    assert title_placeholder["text_defaults"]["suggested_font_size_pt"] == 24.0
    assert title_placeholder["text_defaults"]["suggested_font_family"] == "DB Regular"
    assert title_placeholder["estimated_text_capacity"]["max_lines"] == 2
    assert title_placeholder["estimated_text_capacity"]["source"] == "explicit_guidance"
    assert title_placeholder["estimated_text_capacity"]["confidence"] == "high"
    assert title_placeholder["estimated_text_capacity"]["font_size_pt"] == 24.0

    subtitle_placeholder = next(
        placeholder for placeholder in placeholders if placeholder["logical_name"] == "subtitle"
    )
    assert subtitle_placeholder["estimated_text_capacity"]["max_lines"] == 1
    assert subtitle_placeholder["estimated_text_capacity"]["source"] == "inferred"
    assert subtitle_placeholder["estimated_text_capacity"]["confidence"] == "medium"

    content_layout_payload = _invoke_json(
        ["placeholders", "list", "1-title-and-content", "--manifest", str(manifest_dir)]
    )
    content_placeholder = next(
        placeholder
        for placeholder in content_layout_payload["result"]["placeholders"]
        if placeholder["logical_name"] == "content_1"
    )
    assert content_placeholder["estimated_text_capacity"]["max_lines"] == 20
    assert content_placeholder["estimated_text_capacity"]["source"] == "inferred"
    assert content_placeholder["estimated_text_capacity"]["font_size_pt"] == 14.0

    theme_payload = _invoke_json(["theme", "show", "--manifest", str(manifest_dir)])
    assert "fonts" in theme_payload["result"]
    assert theme_payload["result"]["fonts"]["major"] == "DanskeRegular"
    assert theme_payload["result"]["fonts"]["minor"] == "DanskeRegular"

    assets_payload = _invoke_json(["assets", "list", "--manifest", str(manifest_dir)])
    assert assets_payload["result"]["count"] >= 1

    doctor_payload = _invoke_json(["doctor", "--manifest", str(manifest_dir)])
    assert doctor_payload["result"]["status"] in {"ok", "warn"}


def test_slide_create_and_validate_round_trip(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    out_file = tmp_path / "slide.pptx"
    _init_manifest(template_path, manifest_dir)

    payload = _invoke_json(
        [
            "slide",
            "create",
            "--manifest",
            str(manifest_dir),
            "--layout",
            "title-only",
            "--set",
            "title=Quarterly Update",
            "--set",
            "subtitle=March 2026",
            "--out",
            str(out_file),
        ]
    )

    assert payload["result"]["summary"]["slides"] == 1
    assert out_file.exists()

    validation = _invoke_json(
        [
            "validate",
            "--manifest",
            str(manifest_dir),
            "--deck",
            str(out_file),
        ]
    )
    assert validation["result"]["ok"] is True
    assert validation["result"]["checked_slides"] == 1

    generated = Presentation(str(out_file))
    title_shape: Any = next(
        shape for shape in generated.slides[0].placeholders if shape.placeholder_format.idx == 0
    )
    template = Presentation(str(template_path))
    template_layout = next(
        layout for layout in template.slide_layouts if layout.name == "Title Only"
    )
    template_title_shape: Any = next(
        shape for shape in template_layout.placeholders if shape.placeholder_format.idx == 0
    )
    first_run = title_shape.text_frame.paragraphs[0].runs[0]
    assert first_run.font.size is None
    assert title_shape.text_frame.vertical_anchor == template_title_shape.text_frame.vertical_anchor


def test_slide_create_picture_placeholder_uses_fit_by_default(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    out_file = tmp_path / "picture-slide.pptx"
    image_path = tmp_path / "wide.png"
    _write_png(image_path, width=500, height=100)
    _init_manifest(template_path, manifest_dir)

    payload = _invoke_json(
        [
            "slide",
            "create",
            "--manifest",
            str(manifest_dir),
            "--layout",
            "3-front-page-title-and-picture",
            "--set",
            "title=Workflow",
            "--set",
            f"picture=@{image_path}",
            "--out",
            str(out_file),
        ]
    )

    assert payload["result"]["summary"]["slides"] == 1
    generated = Presentation(str(out_file))
    placeholder_order = [
        shape.placeholder_format.idx
        for shape in generated.slides[0].shapes
        if getattr(shape, "is_placeholder", False)
    ]
    picture_shape = _placeholder_by_idx(generated.slides[0], 14)
    assert picture_shape.crop_top < 0
    assert picture_shape.crop_bottom < 0
    assert picture_shape.crop_left == 0.0
    assert picture_shape.crop_right == 0.0
    assert placeholder_order.index(14) < placeholder_order.index(13)
    assert 0 not in placeholder_order
    assert 19 not in placeholder_order
    assert 17 in placeholder_order
    assert 18 in placeholder_order

    validation = _invoke_json(
        [
            "validate",
            "--manifest",
            str(manifest_dir),
            "--deck",
            str(out_file),
        ]
    )
    assert validation["result"]["ok"] is True


def test_deck_build_supports_structured_image_table_and_chart_content(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    deck_spec = tmp_path / "structured-content.yaml"
    out_file = tmp_path / "structured-content.pptx"
    image_path = tmp_path / "wide.png"
    _write_png(image_path, width=500, height=100)
    _init_manifest(template_path, manifest_dir)

    deck_spec.write_text(
        yaml.safe_dump(
            {
                "metadata": {"title": "Structured content"},
                "slides": [
                    {
                        "layout": "1-title-and-content",
                        "content": {
                            "title": "Image",
                            "content_1": {
                                "kind": "image",
                                "path": str(image_path),
                                "image_fit": "cover",
                            },
                        },
                    },
                    {
                        "layout": "1-title-and-content",
                        "content": {
                            "title": "Table",
                            "content_1": {
                                "kind": "table",
                                "columns": ["Owner", "Status"],
                                "rows": [["Platform", "Active"], ["Sales", "Planned"]],
                            },
                        },
                    },
                    {
                        "layout": "1-title-and-content",
                        "content": {
                            "title": "Chart",
                            "content_1": {
                                "kind": "chart",
                                "chart_type": "column_clustered",
                                "categories": ["Q1", "Q2", "Q3"],
                                "series": [
                                    {"name": "Revenue", "values": [12, 15, 18]},
                                ],
                            },
                        },
                    },
                ],
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    payload = _invoke_json(
        [
            "deck",
            "build",
            "--manifest",
            str(manifest_dir),
            "--spec",
            str(deck_spec),
            "--out",
            str(out_file),
        ]
    )

    assert payload["result"]["summary"]["slides"] == 3

    validation = _invoke_json(
        [
            "validate",
            "--manifest",
            str(manifest_dir),
            "--deck",
            str(out_file),
        ]
    )
    assert validation["result"]["ok"] is True

    generated = Presentation(str(out_file))
    image_shape = _placeholder_by_idx(generated.slides[0], 1)
    assert image_shape.crop_left > 0
    assert image_shape.crop_right > 0
    table_shape = _placeholder_by_idx(generated.slides[1], 1)
    assert table_shape.has_table is True
    chart_shape = _placeholder_by_idx(generated.slides[2], 1)
    assert chart_shape.has_chart is True


def test_deck_build_auto_detects_markdown_bullets_and_preserves_text_opt_out(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    deck_spec = tmp_path / "bullets.yaml"
    out_file = tmp_path / "bullets.pptx"
    _init_manifest(template_path, manifest_dir)

    deck_spec.write_text(
        yaml.safe_dump(
            {
                "slides": [
                    {
                        "layout": "1-title-and-content",
                        "content": {
                            "title": "Bullets",
                            "content_1": "- First point\n  - Nested point\n\nPlain paragraph",
                        },
                    },
                    {
                        "layout": "1-title-and-content",
                        "content": {
                            "title": "Literal dash",
                            "content_1": {"kind": "text", "value": "- literal dash"},
                        },
                    },
                ]
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    payload = _invoke_json(
        [
            "deck",
            "build",
            "--manifest",
            str(manifest_dir),
            "--spec",
            str(deck_spec),
            "--out",
            str(out_file),
        ]
    )

    assert payload["result"]["summary"]["slides"] == 2

    generated = Presentation(str(out_file))
    bullet_shape = _placeholder_by_idx(generated.slides[0], 1)
    paragraphs = bullet_shape.text_frame.paragraphs
    assert [paragraph.text for paragraph in paragraphs] == [
        "First point",
        "Nested point",
        "Plain paragraph",
    ]
    assert paragraphs[0].level == 1
    assert paragraphs[1].level == 2

    literal_shape = _placeholder_by_idx(generated.slides[1], 1)
    assert literal_shape.text_frame.paragraphs[0].text == "- literal dash"


def test_deck_build_preserves_blank_paragraphs_in_plain_multiline_three_column_layout(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    deck_spec = tmp_path / "blank-lines-three-columns.yaml"
    out_file = tmp_path / "blank-lines-three-columns.pptx"
    _init_manifest(template_path, manifest_dir)

    deck_spec.write_text(
        yaml.safe_dump(
            {
                "slides": [
                    {
                        "layout": "1-2-three-contents",
                        "content": {
                            "title": "Plain multiline",
                            "content_1": "First paragraph.\n\nSecond paragraph.",
                            "content_2": "Alpha\n\nBeta",
                            "content_3": "One\n\nTwo",
                        },
                    }
                ]
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    payload = _invoke_json(
        [
            "deck",
            "build",
            "--manifest",
            str(manifest_dir),
            "--spec",
            str(deck_spec),
            "--out",
            str(out_file),
        ]
    )

    assert payload["result"]["summary"]["slides"] == 1

    generated = Presentation(str(out_file))
    left_column = _placeholder_by_idx(generated.slides[0], 1)
    middle_column = _placeholder_by_idx(generated.slides[0], 16)
    right_column = _placeholder_by_idx(generated.slides[0], 17)
    assert [paragraph.text for paragraph in left_column.text_frame.paragraphs] == [
        "First paragraph.",
        "",
        "Second paragraph.",
    ]
    assert [paragraph.text for paragraph in middle_column.text_frame.paragraphs] == [
        "Alpha",
        "",
        "Beta",
    ]
    assert [paragraph.text for paragraph in right_column.text_frame.paragraphs] == [
        "One",
        "",
        "Two",
    ]


def test_deck_build_parses_markdown_inline_formatting_and_ordered_lists(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    deck_spec = tmp_path / "markdown-rich.yaml"
    out_file = tmp_path / "markdown-rich.pptx"
    _init_manifest(template_path, manifest_dir)

    deck_spec.write_text(
        yaml.safe_dump(
            {
                "slides": [
                    {
                        "layout": "1-title-and-content",
                        "content": {
                            "title": "Markdown rich",
                            "content_1": {
                                "kind": "markdown-text",
                                "value": (
                                    "# Heading\n\n"
                                    "Some **bold** and *italic* text.\n\n"
                                    "1. First item\n"
                                    "2. Second item"
                                ),
                            },
                        },
                    }
                ]
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    payload = _invoke_json(
        [
            "deck",
            "build",
            "--manifest",
            str(manifest_dir),
            "--spec",
            str(deck_spec),
            "--out",
            str(out_file),
        ]
    )

    assert payload["result"]["summary"]["slides"] == 1

    generated = Presentation(str(out_file))
    content_shape = _placeholder_by_idx(generated.slides[0], 1)
    paragraphs = content_shape.text_frame.paragraphs
    assert [paragraph.text for paragraph in paragraphs] == [
        "Heading",
        "Some bold and italic text.",
        "1. First item",
        "2. Second item",
    ]

    rich_runs = paragraphs[1].runs
    assert [run.text for run in rich_runs] == ["Some ", "bold", " and ", "italic", " text."]
    assert rich_runs[1].font.bold is True
    assert rich_runs[3].font.italic is True
    assert paragraphs[2].level == 1
    assert paragraphs[3].level == 1


def test_deck_build_applies_markdown_spacing_and_heading_emphasis(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    deck_spec = tmp_path / "markdown-spacing.yaml"
    out_file = tmp_path / "markdown-spacing.pptx"
    _init_manifest(template_path, manifest_dir)

    deck_spec.write_text(
        yaml.safe_dump(
            {
                "slides": [
                    {
                        "layout": "1-title-and-content",
                        "content": {
                            "title": "Markdown spacing",
                            "content_1": {
                                "kind": "markdown-text",
                                "value": (
                                    "# Heading\n\n"
                                    "Paragraph body.\n\n"
                                    "- Bullet one\n"
                                    "- Bullet two\n\n"
                                    "Body after list."
                                ),
                            },
                        },
                    }
                ]
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    payload = _invoke_json(
        [
            "deck",
            "build",
            "--manifest",
            str(manifest_dir),
            "--spec",
            str(deck_spec),
            "--out",
            str(out_file),
        ]
    )

    assert payload["result"]["summary"]["slides"] == 1

    generated = Presentation(str(out_file))
    content_shape = _placeholder_by_idx(generated.slides[0], 1)
    paragraphs = content_shape.text_frame.paragraphs

    assert paragraphs[0].runs[0].font.bold is True
    assert paragraphs[0].space_after is not None
    assert round(paragraphs[0].space_after.pt, 1) == 12.0
    assert paragraphs[1].space_after is not None
    assert round(paragraphs[1].space_after.pt, 1) == 6.0
    assert paragraphs[2].space_before is not None
    assert round(paragraphs[2].space_before.pt, 1) == 6.0
    assert paragraphs[2].space_after is not None
    assert round(paragraphs[2].space_after.pt, 1) == 2.0
    assert paragraphs[4].space_before is not None
    assert round(paragraphs[4].space_before.pt, 1) == 6.0


def test_annotations_can_override_supported_content_types(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "wide.png"
    _write_png(image_path, width=500, height=100)
    _init_manifest(template_path, manifest_dir)

    annotations_path = manifest_dir / "annotations.yaml"
    annotations = yaml.safe_load(annotations_path.read_text(encoding="utf-8"))
    layout_annotation = next(
        item
        for item in annotations["layouts"]
        if item["layout_id"] == "3-front-page-title-and-picture"
    )
    layout_annotation["aliases"] = ["hero-image"]
    layout_annotation["placeholder_overrides"] = [
        {
            "logical_name": "picture",
            "supported_content_types": ["markdown-text"],
        }
    ]
    annotations_path.write_text(
        yaml.safe_dump(annotations, sort_keys=False),
        encoding="utf-8",
    )

    placeholders_payload = _invoke_json(
        [
            "placeholders",
            "list",
            "hero-image",
            "--manifest",
            str(manifest_dir),
        ]
    )
    picture_placeholder = next(
        item
        for item in placeholders_payload["result"]["placeholders"]
        if item["logical_name"] == "picture"
    )
    assert picture_placeholder["supported_content_types"] == ["markdown-text"]

    result = runner.invoke(
        app,
        [
            "slide",
            "create",
            "--manifest",
            str(manifest_dir),
            "--layout",
            "hero-image",
            "--set",
            "title=Workflow",
            "--set",
            f"picture=@{image_path}",
            "--out",
            str(tmp_path / "hero-image.pptx"),
            "--format",
            "json",
        ],
    )

    assert result.exit_code == 10
    payload = json.loads(result.stdout)
    assert payload["errors"][0]["code"] == "ERR_VALIDATION_CONTENT_TYPE"


def test_slide_create_requires_overwrite_for_existing_output(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    out_file = tmp_path / "slide.pptx"
    _init_manifest(template_path, manifest_dir)

    _invoke_json(
        [
            "slide",
            "create",
            "--manifest",
            str(manifest_dir),
            "--layout",
            "title-only",
            "--set",
            "title=Quarterly Update",
            "--out",
            str(out_file),
        ]
    )

    result = runner.invoke(
        app,
        [
            "slide",
            "create",
            "--manifest",
            str(manifest_dir),
            "--layout",
            "title-only",
            "--set",
            "title=Replacement",
            "--out",
            str(out_file),
            "--format",
            "json",
        ],
    )

    assert result.exit_code == 40
    payload = json.loads(result.stdout)
    assert payload["errors"][0]["code"] == "ERR_CONFLICT_OUTPUT_EXISTS"

    overwrite_payload = _invoke_json(
        [
            "slide",
            "create",
            "--manifest",
            str(manifest_dir),
            "--layout",
            "title-only",
            "--set",
            "title=Replacement",
            "--out",
            str(out_file),
            "--overwrite",
        ]
    )
    assert overwrite_payload["result"]["overwrite"] is True


def test_slide_create_fails_for_unknown_placeholder(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    out_file = tmp_path / "slide.pptx"
    _init_manifest(template_path, manifest_dir)

    result = runner.invoke(
        app,
        [
            "slide",
            "create",
            "--manifest",
            str(manifest_dir),
            "--layout",
            "title-only",
            "--set",
            "nope=value",
            "--out",
            str(out_file),
            "--format",
            "json",
        ],
    )

    assert result.exit_code == 10
    payload = json.loads(result.stdout)
    assert payload["errors"][0]["code"] == "ERR_VALIDATION_PLACEHOLDER_UNKNOWN"


def test_deck_build_failure_identifies_slide_and_placeholder_context(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    deck_spec = tmp_path / "missing-image.yaml"
    out_file = tmp_path / "missing-image.pptx"
    missing_image = tmp_path / "missing.png"
    _init_manifest(template_path, manifest_dir)

    deck_spec.write_text(
        yaml.safe_dump(
            {
                "slides": [
                    {
                        "layout": "title-only",
                        "content": {"title": "First slide"},
                    },
                    {
                        "layout": "1-2-three-contents",
                        "content": {
                            "title": "Broken slide",
                            "content_2": {
                                "kind": "image",
                                "path": str(missing_image),
                            },
                        },
                    },
                ]
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    result = runner.invoke(
        app,
        [
            "deck",
            "build",
            "--manifest",
            str(manifest_dir),
            "--spec",
            str(deck_spec),
            "--out",
            str(out_file),
            "--format",
            "json",
        ],
    )

    assert result.exit_code == 50
    payload = json.loads(result.stdout)
    assert payload["errors"][0]["code"] == "ERR_IO_NOT_FOUND"
    assert "Slide 2 (1-2-three-contents):" in payload["errors"][0]["message"]
    assert "Placeholder content_2:" in payload["errors"][0]["message"]
    assert "Image not found:" in payload["errors"][0]["message"]


def test_deck_build_validate_schema_and_diff(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    deck_spec = tmp_path / "deck.yaml"
    out_file = tmp_path / "deck.pptx"
    _init_manifest(template_path, manifest_dir)

    deck_spec.write_text(
        yaml.safe_dump(
            {
                "metadata": {"title": "Operating Model", "author": "Thomas"},
                "slides": [
                    {"layout": "title-only", "content": {"title": "Operating Model"}},
                    {
                        "layout": "1-title-and-content",
                        "content": {
                            "title": "Core idea",
                            "content_1": "Stay inside the template contract.",
                            "subtitle": "Programmatic deck build",
                        },
                    },
                ],
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    build_payload = _invoke_json(
        [
            "deck",
            "build",
            "--manifest",
            str(manifest_dir),
            "--spec",
            str(deck_spec),
            "--out",
            str(out_file),
        ]
    )
    assert build_payload["result"]["summary"]["slides"] == 2
    assert out_file.exists()

    schema_payload = _invoke_json(["manifest", "schema"])
    assert "template" in schema_payload["result"]["properties"]
    assert "$defs" in schema_payload["result"]
    assert (
        "estimated_text_capacity"
        in schema_payload["result"]["$defs"]["PlaceholderContract"]["properties"]
    )

    diff_payload = _invoke_json(["manifest", "diff", str(manifest_dir), str(manifest_dir)])
    assert diff_payload["result"]["breaking_changes"] == []
    assert "layouts" in diff_payload["result"]["unchanged"]

    modified_manifest_dir = tmp_path / "manifest-copy"
    shutil.copytree(manifest_dir, modified_manifest_dir)
    manifest_path = modified_manifest_dir / "manifest.yaml"
    manifest_payload = yaml.safe_load(manifest_path.read_text(encoding="utf-8"))
    target_placeholder = next(
        placeholder
        for layout in manifest_payload["layouts"]
        for placeholder in layout["placeholders"]
        if placeholder["placeholder_type"] == "picture" and placeholder["logical_name"] == "picture"
    )
    target_placeholder["supported_content_types"] = ["image", "markdown-text"]
    manifest_path.write_text(
        yaml.safe_dump(manifest_payload, sort_keys=False),
        encoding="utf-8",
    )

    changed_diff_payload = _invoke_json(
        ["manifest", "diff", str(manifest_dir), str(modified_manifest_dir)]
    )
    assert any(
        change["type"] == "placeholder.content_types_changed"
        for change in changed_diff_payload["result"]["breaking_changes"]
    )


def test_wrapper_generate_dry_run_and_apply(
    template_path: Path,
    manifest_dir: Path,
    tmp_path: Path,
) -> None:
    wrapper_dir = tmp_path / "wrapper"
    _init_manifest(template_path, manifest_dir)

    dry_run_payload = _invoke_json(
        [
            "wrapper",
            "generate",
            "--manifest",
            str(manifest_dir),
            "--out",
            str(wrapper_dir),
            "--dry-run",
        ]
    )
    assert dry_run_payload["result"]["dry_run"] is True
    assert dry_run_payload["result"]["summary"]["artifacts"] >= 4

    apply_payload = _invoke_json(
        [
            "wrapper",
            "generate",
            "--manifest",
            str(manifest_dir),
            "--out",
            str(wrapper_dir),
        ]
    )
    assert apply_payload["result"]["dry_run"] is False
    assert (wrapper_dir / "pyproject.toml").exists()
