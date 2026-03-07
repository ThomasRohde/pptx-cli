from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from pptx_cli.models.manifest import (
    ManifestDiffResult,
    ManifestDocument,
    ValidationIssue,
    ValidationResult,
)


class ValidationError(ValueError):
    code: str

    def __init__(self, code: str, message: str) -> None:
        super().__init__(message)
        self.code = code


def validate_deck(
    manifest_dir: Path,
    manifest: ManifestDocument,
    deck_path: Path,
    *,
    strict: bool,
) -> ValidationResult:
    if not deck_path.exists():
        raise ValidationError("ERR_IO_NOT_FOUND", f"Deck not found: {deck_path}")

    prs = Presentation(str(deck_path))
    issues: list[ValidationIssue] = []
    expected_size = manifest.presentation.get("page_size", {})
    actual_width = int(prs.slide_width or 0)
    actual_height = int(prs.slide_height or 0)
    if actual_width != expected_size.get("width_emu") or actual_height != expected_size.get(
        "height_emu"
    ):
        issues.append(
            ValidationIssue(
                code="ERR_VALIDATION_PAGE_SIZE",
                severity="error",
                message="Deck page size does not match the manifest",
                details={
                    "expected": expected_size,
                    "actual": {"width_emu": actual_width, "height_emu": actual_height},
                },
            )
        )

    layouts_by_source_name = {layout.source_layout_name: layout for layout in manifest.layouts}
    checked_layouts: set[str] = set()
    for slide_index, slide in enumerate(prs.slides, start=1):
        source_layout = getattr(slide, "slide_layout", None)
        layout_name = getattr(source_layout, "name", None)
        if layout_name not in layouts_by_source_name:
            issues.append(
                ValidationIssue(
                    code="ERR_VALIDATION_LAYOUT_UNKNOWN",
                    severity="error",
                    message="Slide uses a layout that is not present in the manifest",
                    details={"slide_index": slide_index, "layout_name": layout_name},
                )
            )
            continue

        layout = layouts_by_source_name[layout_name]
        checked_layouts.add(layout.id)
        expected_idxs = {placeholder.placeholder_idx for placeholder in layout.placeholders}
        actual_idxs = {shape.placeholder_format.idx for shape in slide.placeholders}
        missing_idxs = sorted(expected_idxs - actual_idxs)
        if missing_idxs:
            issues.append(
                ValidationIssue(
                    code="ERR_VALIDATION_PLACEHOLDER_MISSING",
                    severity="error",
                    message="Slide is missing expected placeholders for its layout",
                    details={
                        "slide_index": slide_index,
                        "layout_id": layout.id,
                        "placeholder_idxs": missing_idxs,
                    },
                )
            )

    if manifest.compatibility_report.findings:
        for finding in manifest.compatibility_report.findings:
            if finding.severity == "warning":
                issues.append(
                    ValidationIssue(
                        code=finding.code,
                        severity="error" if strict else "warning",
                        message=f"Manifest compatibility warning: {finding.message}",
                        details=finding.details,
                    )
                )

    has_errors = any(issue.severity == "error" for issue in issues)
    return ValidationResult(
        manifest_path=str(manifest_dir),
        deck_path=str(deck_path),
        ok=not has_errors,
        issues=issues,
        checked_slides=len(prs.slides),
        checked_layouts=len(checked_layouts),
    )


def diff_manifests(left: ManifestDocument, right: ManifestDocument) -> ManifestDiffResult:
    result = ManifestDiffResult()
    left_layouts = {layout.id: layout for layout in left.layouts}
    right_layouts = {layout.id: layout for layout in right.layouts}

    removed_layouts = sorted(set(left_layouts) - set(right_layouts))
    added_layouts = sorted(set(right_layouts) - set(left_layouts))

    for layout_id in removed_layouts:
        result.breaking_changes.append({"type": "layout.removed", "layout_id": layout_id})
    for layout_id in added_layouts:
        result.additive_changes.append({"type": "layout.added", "layout_id": layout_id})

    shared_layouts = sorted(set(left_layouts) & set(right_layouts))
    for layout_id in shared_layouts:
        left_layout = left_layouts[layout_id]
        right_layout = right_layouts[layout_id]
        if left_layout.source_layout_name != right_layout.source_layout_name:
            result.breaking_changes.append(
                {
                    "type": "layout.renamed",
                    "layout_id": layout_id,
                    "before": left_layout.source_layout_name,
                    "after": right_layout.source_layout_name,
                }
            )
        left_placeholders = {
            placeholder.logical_name: placeholder for placeholder in left_layout.placeholders
        }
        right_placeholders = {
            placeholder.logical_name: placeholder for placeholder in right_layout.placeholders
        }
        removed_placeholders = sorted(set(left_placeholders) - set(right_placeholders))
        added_placeholders = sorted(set(right_placeholders) - set(left_placeholders))
        for placeholder_name in removed_placeholders:
            result.breaking_changes.append(
                {
                    "type": "placeholder.removed",
                    "layout_id": layout_id,
                    "placeholder": placeholder_name,
                }
            )
        for placeholder_name in added_placeholders:
            result.additive_changes.append(
                {
                    "type": "placeholder.added",
                    "layout_id": layout_id,
                    "placeholder": placeholder_name,
                }
            )
        for placeholder_name in sorted(set(left_placeholders) & set(right_placeholders)):
            left_placeholder = left_placeholders[placeholder_name]
            right_placeholder = right_placeholders[placeholder_name]
            left_geometry = (
                left_placeholder.left_emu,
                left_placeholder.top_emu,
                left_placeholder.width_emu,
                left_placeholder.height_emu,
            )
            right_geometry = (
                right_placeholder.left_emu,
                right_placeholder.top_emu,
                right_placeholder.width_emu,
                right_placeholder.height_emu,
            )
            if left_geometry != right_geometry:
                result.breaking_changes.append(
                    {
                        "type": "placeholder.geometry_changed",
                        "layout_id": layout_id,
                        "placeholder": placeholder_name,
                        "before": left_geometry,
                        "after": right_geometry,
                    }
                )
            if (
                left_placeholder.supported_content_types
                != right_placeholder.supported_content_types
            ):
                result.breaking_changes.append(
                    {
                        "type": "placeholder.content_types_changed",
                        "layout_id": layout_id,
                        "placeholder": placeholder_name,
                        "before": left_placeholder.supported_content_types,
                        "after": right_placeholder.supported_content_types,
                    }
                )

    left_theme = left.presentation.get("theme", {})
    right_theme = right.presentation.get("theme", {})
    if left_theme != right_theme:
        result.breaking_changes.append(
            {
                "type": "theme.changed",
                "before": left_theme,
                "after": right_theme,
            }
        )

    if not result.breaking_changes and not result.additive_changes:
        result.unchanged.append("layouts")
        result.unchanged.append("theme")

    return result
