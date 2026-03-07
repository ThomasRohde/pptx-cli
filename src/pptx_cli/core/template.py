from __future__ import annotations

import hashlib
import math
import re
import shutil
import zipfile
from collections import Counter
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

import lxml.etree as etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from pptx_cli.core.ids import slugify, uniquify
from pptx_cli.core.io import atomic_write_bytes, ensure_directory
from pptx_cli.core.manifest_store import template_copy_path
from pptx_cli.models.manifest import (
    AnnotationsDocument,
    AssetRef,
    CompatibilityFinding,
    CompatibilityReport,
    InitReport,
    LayoutAnnotation,
    LayoutContract,
    ManifestDocument,
    MasterContract,
    PlaceholderContract,
    ProtectedElement,
    TemplateInfo,
    TextCapacityGuidance,
    ThemeModel,
)

_DRAWINGML_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
_PRESENTATIONML_NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
_PPTX_NS = {**_DRAWINGML_NS, **_PRESENTATIONML_NS}
_PLACEHOLDER_TYPE_NAMES = {item.value: item.name.lower() for item in PP_PLACEHOLDER}
_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg", ".tif", ".tiff"}
_MEDIA_SUFFIXES = {".mp4", ".wmv", ".avi", ".mov", ".mp3", ".wav", ".m4v"}
_EMU_PER_POINT = 12700
_DEFAULT_LINE_HEIGHT_MULTIPLIER = 1.2
_NUMBER_WORDS = {
    "one": 1,
    "two": 2,
    "three": 3,
    "four": 4,
    "five": 5,
    "six": 6,
    "seven": 7,
    "eight": 8,
    "nine": 9,
    "ten": 10,
}


def sha256_bytes(payload: bytes) -> str:
    return f"sha256:{hashlib.sha256(payload).hexdigest()}"


def sha256_file(path: Path) -> str:
    return sha256_bytes(path.read_bytes())


def _placeholder_type_name(value: int) -> str:
    return _PLACEHOLDER_TYPE_NAMES.get(value, f"placeholder_{value}")


def _supports_content_types(placeholder_type: str, shape_name: str) -> list[str]:
    lower_name = shape_name.lower()
    if "logo" in lower_name or "progress bar" in lower_name:
        return []
    if placeholder_type in {
        "pic",
        "picture",
        "bitmap",
        "media_clip",
        "org_chart",
        "clip_art",
        "slide_image",
    }:
        return ["image"]
    if placeholder_type == "chart":
        return ["chart"]
    if placeholder_type == "table":
        return ["table"]
    if placeholder_type in {"body", "object", "content", "text"}:
        return ["text", "markdown-text", "image", "table", "chart"]
    return ["text", "markdown-text"]


def _logical_placeholder_name(source_name: str, placeholder_type: str, placeholder_idx: int) -> str:
    lower_name = source_name.lower()
    if "title" in lower_name and "subtitle" not in lower_name:
        return "title"
    if "subtitle" in lower_name:
        return "subtitle"
    if "date" in lower_name:
        return "date"
    if "source" in lower_name:
        return "source"
    if "description" in lower_name:
        suffix = "".join(ch for ch in source_name if ch.isdigit())
        return f"description_{suffix}" if suffix else "description"
    if "picture" in lower_name:
        return "picture"
    if "content" in lower_name:
        suffix = "".join(ch for ch in source_name if ch.isdigit())
        return f"content_{suffix}" if suffix else "content"
    return f"{slugify(placeholder_type)}_{placeholder_idx}"


def _shape_fingerprint(shape: Any) -> str:
    element = getattr(shape, "_element", None)
    xml = etree.tostring(element) if element is not None else shape.name.encode("utf-8")
    return sha256_bytes(xml)


def _guidance_lines(shape: Any) -> list[str]:
    raw_text = getattr(shape, "text", "")
    normalized = raw_text.replace("\v", "\n")
    return [line.strip() for line in normalized.splitlines() if line.strip()]


def _parse_max_lines(lines: list[str]) -> int | None:
    for line in lines:
        match = re.search(r"max\s+(?P<count>\d+|[A-Za-z]+)\s+lines?", line, re.IGNORECASE)
        if match is None:
            continue
        count_value = match.group("count").lower()
        if count_value.isdigit():
            return int(count_value)
        return _NUMBER_WORDS.get(count_value)
    return None


def _parse_suggested_font_size(text: str) -> float | None:
    match = re.search(r"(?P<size>\d+(?:\.\d+)?)\s*pt", text, re.IGNORECASE)
    if match is None:
        return None
    return float(match.group("size"))


def _parse_suggested_font_family(text: str) -> str | None:
    match = re.search(
        r"font size\s+(?P<family>.+?)\s+\d+(?:\.\d+)?\s*pt",
        text,
        re.IGNORECASE,
    )
    if match is None:
        return None
    return match.group("family").strip()


def _extract_text_defaults(shape: Any) -> dict[str, Any]:
    guidance_lines = _guidance_lines(shape)
    guidance_text = "\n".join(guidance_lines)
    text_frame = getattr(shape, "text_frame", None)
    alignment = None
    paragraph_count = 0
    if text_frame is not None:
        paragraph_count = len(text_frame.paragraphs)
        if text_frame.paragraphs:
            paragraph_alignment = text_frame.paragraphs[0].alignment
            alignment = paragraph_alignment.name.lower() if paragraph_alignment else None

    defaults = {
        "guidance_text": guidance_text or None,
        "guidance_lines": guidance_lines,
        "max_lines": _parse_max_lines(guidance_lines),
        "suggested_font_size_pt": _parse_suggested_font_size(guidance_text),
        "suggested_font_family": _parse_suggested_font_family(guidance_text),
        "paragraph_count": paragraph_count,
        "alignment": alignment,
    }
    return {key: value for key, value in defaults.items() if value not in (None, [], "")}


def _emu_to_points(value: int) -> float:
    return round(value / _EMU_PER_POINT, 2)


def _placeholder_text_style_bucket(placeholder_type: str) -> str:
    if placeholder_type in {"title", "center_title", "vertical_title"}:
        return "title"
    if placeholder_type in {"body", "object", "content", "text", "subtitle"}:
        return "body"
    return "other"


def _extract_master_text_styles(master: Any) -> dict[str, dict[int, float]]:
    styles: dict[str, dict[int, float]] = {}
    for bucket, tag_name in {
        "title": "titleStyle",
        "body": "bodyStyle",
        "other": "otherStyle",
    }.items():
        style_root = master.element.find(f".//p:{tag_name}", namespaces=_PPTX_NS)
        if style_root is None:
            styles[bucket] = {}
            continue

        level_sizes: dict[int, float] = {}
        default_size_pt: float | None = None
        default_paragraph = style_root.find("./a:defPPr", namespaces=_PPTX_NS)
        if default_paragraph is not None:
            default_run = default_paragraph.find("./a:defRPr", namespaces=_PPTX_NS)
            default_size_pt = _font_size_from_xml(default_run)

        for level in range(1, 10):
            paragraph_style = style_root.find(f"./a:lvl{level}pPr", namespaces=_PPTX_NS)
            if paragraph_style is None:
                continue
            default_run = paragraph_style.find("./a:defRPr", namespaces=_PPTX_NS)
            font_size_pt = _font_size_from_xml(default_run)
            if font_size_pt is not None:
                level_sizes[level - 1] = font_size_pt

        if default_size_pt is not None:
            level_sizes.setdefault(0, default_size_pt)
        styles[bucket] = level_sizes

    return styles


def _font_size_from_xml(element: Any) -> float | None:
    if element is None:
        return None
    size = element.get("sz")
    if size is None:
        return None
    try:
        return int(size) / 100
    except ValueError:
        return None


def _resolve_font_size_pt(
    shape: Any,
    placeholder_type: str,
    text_defaults: dict[str, Any],
    master_text_styles: dict[str, dict[int, float]],
) -> tuple[float | None, str | None]:
    suggested_size = text_defaults.get("suggested_font_size_pt")
    if suggested_size is not None:
        return float(suggested_size), "guidance_text"

    text_frame = getattr(shape, "text_frame", None)
    if text_frame is not None:
        for paragraph in text_frame.paragraphs:
            if paragraph.font.size is not None:
                return float(paragraph.font.size.pt), "paragraph_font"
            for run in paragraph.runs:
                if run.font.size is not None:
                    return float(run.font.size.pt), "run_font"

    bucket = _placeholder_text_style_bucket(placeholder_type)
    text_level = 0
    if text_frame is not None and text_frame.paragraphs:
        text_level = int(text_frame.paragraphs[0].level or 0)
    bucket_styles = master_text_styles.get(bucket, {})
    font_size_pt = bucket_styles.get(text_level)
    if font_size_pt is not None:
        return float(font_size_pt), "master_text_style"
    if bucket != "other":
        font_size_pt = master_text_styles.get("other", {}).get(text_level)
        if font_size_pt is not None:
            return float(font_size_pt), "master_text_style"
    return None, None


def _resolve_font_family(
    shape: Any,
    placeholder_type: str,
    text_defaults: dict[str, Any],
    theme: ThemeModel,
) -> tuple[str | None, str | None]:
    suggested_family = text_defaults.get("suggested_font_family")
    if suggested_family is not None:
        return str(suggested_family), "guidance_text"

    text_frame = getattr(shape, "text_frame", None)
    if text_frame is not None:
        for paragraph in text_frame.paragraphs:
            if paragraph.font.name:
                return str(paragraph.font.name), "paragraph_font"
            for run in paragraph.runs:
                if run.font.name:
                    return str(run.font.name), "run_font"

    bucket = _placeholder_text_style_bucket(placeholder_type)
    if bucket == "title":
        theme_family = theme.fonts.get("major") or theme.fonts.get("major_latin")
    else:
        theme_family = theme.fonts.get("minor") or theme.fonts.get("minor_latin")
    if theme_family is None:
        theme_family = theme.fonts.get("major") or theme.fonts.get("minor")
    if theme_family is None:
        return None, None
    return theme_family, "theme"


def _estimate_text_capacity(
    shape: Any,
    placeholder_type: str,
    supported_content_types: list[str],
    text_defaults: dict[str, Any],
    master_text_styles: dict[str, dict[int, float]],
    theme: ThemeModel,
) -> TextCapacityGuidance | None:
    if "text" not in supported_content_types and "markdown-text" not in supported_content_types:
        return None

    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return None

    font_size_pt, font_size_source = _resolve_font_size_pt(
        shape,
        placeholder_type,
        text_defaults,
        master_text_styles,
    )
    if font_size_pt is None:
        return None

    font_family, _ = _resolve_font_family(shape, placeholder_type, text_defaults, theme)
    margin_top = int(text_frame.margin_top or 0)
    margin_bottom = int(text_frame.margin_bottom or 0)
    usable_height_emu = max(int(shape.height) - margin_top - margin_bottom, 0)
    usable_height_pt = _emu_to_points(usable_height_emu)
    if usable_height_pt <= 0:
        return None

    line_height_pt = round(font_size_pt * _DEFAULT_LINE_HEIGHT_MULTIPLIER, 2)
    explicit_max_lines = text_defaults.get("max_lines")
    if explicit_max_lines is not None:
        return TextCapacityGuidance(
            max_lines=int(explicit_max_lines),
            source="explicit_guidance",
            confidence="high",
            font_size_pt=round(font_size_pt, 2),
            font_family=font_family,
            usable_height_pt=usable_height_pt,
            line_height_pt=line_height_pt,
        )

    inferred_lines = max(1, math.floor(usable_height_pt / line_height_pt))
    confidence = "medium" if font_size_source is not None else "low"
    return TextCapacityGuidance(
        max_lines=inferred_lines,
        source="inferred",
        confidence=confidence,
        font_size_pt=round(font_size_pt, 2),
        font_family=font_family,
        usable_height_pt=usable_height_pt,
        line_height_pt=line_height_pt,
    )


def _extract_theme(zip_file: zipfile.ZipFile) -> ThemeModel:
    theme_candidates = [
        name
        for name in zip_file.namelist()
        if name.startswith("ppt/theme/") and name.endswith(".xml")
    ]
    if not theme_candidates:
        return ThemeModel()

    payload = zip_file.read(theme_candidates[0])
    root = etree.fromstring(payload)
    theme_name = root.get("name")
    color_scheme = root.find(".//a:clrScheme", namespaces=_DRAWINGML_NS)
    font_scheme = root.find(".//a:fontScheme", namespaces=_DRAWINGML_NS)

    colors: dict[str, str] = {}
    if color_scheme is not None:
        for child in color_scheme:
            child_value = None
            if len(child):
                inner = child[0]
                child_value = inner.get("val") or inner.get("lastClr")
            if child_value is not None:
                colors[etree.QName(child).localname] = child_value

    fonts: dict[str, str] = {}
    if font_scheme is not None:
        major_latin = font_scheme.find("./a:majorFont/a:latin", namespaces=_DRAWINGML_NS)
        minor_latin = font_scheme.find("./a:minorFont/a:latin", namespaces=_DRAWINGML_NS)
        major_typeface = major_latin.get("typeface") if major_latin is not None else None
        minor_typeface = minor_latin.get("typeface") if minor_latin is not None else None
        if major_typeface:
            fonts["major"] = major_typeface
            fonts["major_latin"] = major_typeface
            fonts["latin_1"] = major_typeface
        if minor_typeface:
            fonts["minor"] = minor_typeface
            fonts["minor_latin"] = minor_typeface
            fonts["latin_2"] = minor_typeface

    return ThemeModel(name=theme_name, colors=colors, fonts=fonts, effects={})


def _layout_description(layout_name: str) -> str | None:
    lower_name = layout_name.lower()
    if "front page" in lower_name:
        return "Front-page style layout extracted from the source template."
    if "breaker" in lower_name:
        return "Section divider layout extracted from the source template."
    if "agenda" in lower_name:
        return "Agenda layout extracted from the source template."
    if "title" in lower_name and "content" in lower_name:
        return "Content layout extracted from the source template."
    if "blank" in lower_name:
        return "Minimal layout with locked brand elements only."
    return None


def _is_protected_placeholder(shape_name: str) -> bool:
    lower_name = shape_name.lower()
    return "logo" in lower_name or "progress bar" in lower_name


def _build_layouts(
    prs: Any,
    theme: ThemeModel,
) -> tuple[list[MasterContract], list[LayoutContract], list[LayoutAnnotation]]:
    master_ids: list[MasterContract] = []
    layouts: list[LayoutContract] = []
    annotations: list[LayoutAnnotation] = []
    layout_ids_seen: set[str] = set()

    master_id_map: dict[int, str] = {}
    master_style_map: dict[int, dict[str, dict[int, float]]] = {}
    for master_index, master in enumerate(prs.slide_masters):
        master_id = f"master-{master_index + 1}"
        master_id_map[id(master)] = master_id
        master_style_map[id(master)] = _extract_master_text_styles(master)
        master_ids.append(
            MasterContract(
                id=master_id,
                name=f"Slide Master {master_index + 1}",
                layout_ids=[],
            )
        )

    for layout_index, layout in enumerate(prs.slide_layouts):
        layout_id = uniquify(slugify(layout.name), layout_ids_seen)
        master_id = master_id_map.get(id(layout.slide_master), "master-1")
        logical_names_seen: set[str] = set()
        placeholders: list[PlaceholderContract] = []
        protected_elements: list[ProtectedElement] = []

        for shape in layout.shapes:
            if getattr(shape, "is_placeholder", False):
                placeholder_format = shape.placeholder_format
                placeholder_type = _placeholder_type_name(int(placeholder_format.type))
                if _is_protected_placeholder(shape.name):
                    protected_elements.append(
                        ProtectedElement(
                            element_id=f"{layout_id}-protected-{placeholder_format.idx}",
                            element_type=placeholder_type,
                            name=shape.name,
                            left_emu=int(shape.left),
                            top_emu=int(shape.top),
                            width_emu=int(shape.width),
                            height_emu=int(shape.height),
                            fingerprint=_shape_fingerprint(shape),
                        )
                    )
                    continue

                logical_name = uniquify(
                    _logical_placeholder_name(shape.name, placeholder_type, placeholder_format.idx),
                    logical_names_seen,
                )
                supported_content_types = _supports_content_types(placeholder_type, shape.name)
                text_defaults = _extract_text_defaults(shape)
                estimated_text_capacity = _estimate_text_capacity(
                    shape,
                    placeholder_type,
                    supported_content_types,
                    text_defaults,
                    master_style_map.get(id(layout.slide_master), {}),
                    theme,
                )
                placeholders.append(
                    PlaceholderContract(
                        logical_name=logical_name,
                        source_name=shape.name,
                        placeholder_idx=int(placeholder_format.idx),
                        placeholder_type=placeholder_type,
                        guidance_text=("\n".join(_guidance_lines(shape)) or None),
                        guidance_lines=_guidance_lines(shape),
                        supported_content_types=supported_content_types,
                        left_emu=int(shape.left),
                        top_emu=int(shape.top),
                        width_emu=int(shape.width),
                        height_emu=int(shape.height),
                        required=logical_name == "title",
                        text_defaults=text_defaults,
                        estimated_text_capacity=estimated_text_capacity,
                        inheritance_chain=[master_id, layout_id],
                    )
                )
                continue

            protected_elements.append(
                ProtectedElement(
                    element_id=f"{layout_id}-static-{len(protected_elements) + 1}",
                    element_type=str(getattr(shape, "shape_type", "shape")),
                    name=shape.name,
                    left_emu=int(shape.left),
                    top_emu=int(shape.top),
                    width_emu=int(shape.width),
                    height_emu=int(shape.height),
                    fingerprint=_shape_fingerprint(shape),
                )
            )

        layouts.append(
            LayoutContract(
                id=layout_id,
                name=layout.name,
                source_master_id=master_id,
                source_layout_index=layout_index,
                source_layout_name=layout.name,
                description=_layout_description(layout.name),
                preview_path=f"previews/layouts/{layout_id}.png",
                placeholders=placeholders,
                protected_static_elements=protected_elements,
                validation_rules={
                    "required_placeholders": [
                        item.logical_name for item in placeholders if item.required
                    ],
                    "protected_elements_locked": True,
                },
            )
        )
        annotations.append(LayoutAnnotation(layout_id=layout_id))

    for layout in layouts:
        for master in master_ids:
            if layout.source_master_id == master.id:
                master.layout_ids.append(layout.id)
                break

    return master_ids, layouts, annotations


def _copy_template_and_assets(template: Path, output_dir: Path) -> list[AssetRef]:
    assets: list[AssetRef] = []
    template_destination = template_copy_path(output_dir)
    ensure_directory(template_destination.parent)
    atomic_write_bytes(template_destination, template.read_bytes())
    assets.append(
        AssetRef(
            id="source-template",
            kind="template",
            path=str(template_destination.relative_to(output_dir)),
            source_path=str(template),
            sha256=sha256_file(template),
            size_bytes=template.stat().st_size,
        )
    )

    with zipfile.ZipFile(template) as zip_file:
        for part_name in zip_file.namelist():
            if part_name.endswith("/"):
                continue
            payload = zip_file.read(part_name)
            relative_target: Path | None = None
            kind: str | None = None
            if part_name.startswith("ppt/media/"):
                suffix = Path(part_name).suffix.lower()
                kind = "media" if suffix in _MEDIA_SUFFIXES else "image"
                asset_dir = "media" if kind == "media" else "images"
                relative_target = Path("assets") / asset_dir / Path(part_name).name
            elif part_name.startswith("ppt/embeddings/"):
                kind = "embedded"
                relative_target = Path("assets") / "embedded" / Path(part_name).name
            elif part_name.startswith("ppt/theme/"):
                kind = "theme"
                relative_target = Path("assets") / "theme" / Path(part_name).name

            if relative_target is None or kind is None:
                continue

            destination = output_dir / relative_target
            ensure_directory(destination.parent)
            atomic_write_bytes(destination, payload)
            assets.append(
                AssetRef(
                    id=slugify(Path(part_name).stem),
                    kind=kind,  # type: ignore[arg-type]
                    path=str(relative_target).replace("\\", "/"),
                    source_path=part_name,
                    sha256=sha256_bytes(payload),
                    size_bytes=len(payload),
                )
            )

    return assets


def _compatibility_findings(template: Path) -> list[CompatibilityFinding]:
    findings: list[CompatibilityFinding] = []
    with zipfile.ZipFile(template) as zip_file:
        names = zip_file.namelist()
        if any(name.startswith("ppt/embeddings/") for name in names):
            findings.append(
                CompatibilityFinding(
                    code="WARN_UNSUPPORTED_EMBEDDINGS",
                    severity="warning",
                    message=(
                        "Embedded OLE objects were found; advanced embedded-object "
                        "fidelity is best-effort in v1."
                    ),
                )
            )
        if any(Path(name).suffix.lower() in _MEDIA_SUFFIXES for name in names):
            findings.append(
                CompatibilityFinding(
                    code="WARN_UNSUPPORTED_MEDIA",
                    severity="warning",
                    message=(
                        "Audio or video media were found; advanced media fidelity "
                        "is out of scope for v1."
                    ),
                )
            )

        xml_tag_counts = Counter[str]()
        for name in names:
            if not name.endswith(".xml"):
                continue
            payload = zip_file.read(name)
            if b"<p:transition" in payload:
                xml_tag_counts["transition"] += 1
            if b"<p:anim" in payload or b"<p:animClr" in payload:
                xml_tag_counts["animation"] += 1

        if xml_tag_counts["transition"]:
            findings.append(
                CompatibilityFinding(
                    code="WARN_UNSUPPORTED_TRANSITIONS",
                    severity="warning",
                    message=(
                        "Slide transitions were detected; transitions are not "
                        "preserved by the v1 fidelity contract."
                    ),
                    details={"count": xml_tag_counts["transition"]},
                )
            )
        if xml_tag_counts["animation"]:
            findings.append(
                CompatibilityFinding(
                    code="WARN_UNSUPPORTED_ANIMATIONS",
                    severity="warning",
                    message=(
                        "Animations were detected; animations are not preserved by "
                        "the v1 fidelity contract."
                    ),
                    details={"count": xml_tag_counts["animation"]},
                )
            )

    if not findings:
        findings.append(
            CompatibilityFinding(
                code="INFO_TEMPLATE_ANALYZED",
                severity="info",
                message=(
                    "Template analysis completed with no known unsupported constructs detected."
                ),
            )
        )
    return findings


def _part_fingerprints(template: Path) -> dict[str, str]:
    with zipfile.ZipFile(template) as zip_file:
        interesting = [
            name
            for name in zip_file.namelist()
            if name == "ppt/presentation.xml"
            or name.startswith("ppt/theme/")
            or name.startswith("ppt/slideMasters/")
            or name.startswith("ppt/slideLayouts/")
        ]
        return {name: sha256_bytes(zip_file.read(name)) for name in sorted(interesting)}


def _presentation_metadata(prs: Any, theme: ThemeModel) -> dict[str, Any]:
    slide_count = len(prs.slides)
    return {
        "page_size": {
            "width_emu": int(prs.slide_width),
            "height_emu": int(prs.slide_height),
        },
        "slide_count": slide_count,
        "theme": theme.model_dump(mode="json", exclude_none=True),
    }


def build_manifest_package(
    template: Path,
    output_dir: Path,
) -> tuple[ManifestDocument, AnnotationsDocument, InitReport]:
    prs = Presentation(str(template))
    with zipfile.ZipFile(template) as zip_file:
        theme = _extract_theme(zip_file)
    masters, layouts, annotations = _build_layouts(prs, theme)
    assets = _copy_template_and_assets(template, output_dir)
    findings = _compatibility_findings(template)
    has_errors = any(item.severity == "error" for item in findings)
    has_warnings = any(item.severity == "warning" for item in findings)
    if has_errors:
        compatibility_status = "error"
    elif has_warnings:
        compatibility_status = "warn"
    else:
        compatibility_status = "ok"
    manifest = ManifestDocument(
        template=TemplateInfo(
            name=template.stem,
            source_file=template.name,
            source_hash=sha256_file(template),
            extracted_at=datetime.now(UTC),
            stored_template_path="assets/source-template.pptx",
        ),
        presentation=_presentation_metadata(prs, theme),
        masters=masters,
        layouts=layouts,
        assets=assets,
        rules={
            "default_policy_mode": "warn",
            "strict_supported": True,
            "supported_placeholder_content_types": [
                "text",
                "image",
                "table",
                "chart",
                "markdown-text",
            ],
            "default_image_fit": "fit",
            "safe_writes": True,
        },
        capabilities={
            "preview_rendering": False,
            "wrapper_generation": True,
            "deck_build": True,
            "slide_create": True,
            "validation": True,
            "manifest_diff": True,
        },
        compatibility_report=CompatibilityReport(status=compatibility_status, findings=findings),
        fingerprints=_part_fingerprints(template),
    )
    annotations_document = AnnotationsDocument(layouts=annotations)
    placeholder_count = sum(len(layout.placeholders) for layout in layouts)
    init_report = InitReport(
        template=str(template),
        output_dir=str(output_dir),
        manifest_path=str(output_dir / "manifest.yaml"),
        findings=findings,
        assets_copied=len(assets),
        layout_count=len(layouts),
        placeholder_count=placeholder_count,
    )
    return manifest, annotations_document, init_report


def plan_manifest_writes(template: Path, output_dir: Path) -> list[dict[str, str]]:
    targets = [
        output_dir / "manifest.yaml",
        output_dir / "manifest.schema.json",
        output_dir / "annotations.yaml",
        output_dir / "reports/init-report.json",
        output_dir / "assets/source-template.pptx",
        output_dir / "previews/layouts/.keep",
        output_dir / "fingerprints/parts.json",
    ]
    changes: list[dict[str, str]] = []
    for target in targets:
        operation = "replace" if target.exists() else "create"
        changes.append({"target": str(target), "operation": operation})
    return changes


def ensure_manifest_directories(output_dir: Path) -> None:
    for relative_path in [
        Path("assets/images"),
        Path("assets/media"),
        Path("assets/embedded"),
        Path("assets/theme"),
        Path("previews/layouts"),
        Path("fingerprints"),
        Path("reports"),
    ]:
        ensure_directory(output_dir / relative_path)


def write_fingerprints(output_dir: Path, manifest: ManifestDocument) -> None:
    from pptx_cli.core.io import write_json

    write_json(output_dir / "fingerprints/parts.json", manifest.fingerprints)
    keep_file = output_dir / "previews/layouts/.keep"
    if not keep_file.exists():
        keep_file.write_text("preview rendering deferred in v1\n", encoding="utf-8")


def copy_output_tree(source: Path, destination: Path) -> None:
    shutil.copytree(source, destination, dirs_exist_ok=True)
