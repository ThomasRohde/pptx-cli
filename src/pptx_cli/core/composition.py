from __future__ import annotations

import json
import os
import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.oxml.shapes.graphfrm import CT_GraphicalObjectFrame
from pptx.oxml.shapes.picture import CT_Picture
from pptx.oxml.xmlchemy import OxmlElement
from pptx.shapes.placeholder import PlaceholderGraphicFrame, PlaceholderPicture
from pptx.util import Emu, Pt

from pptx_cli.core.io import load_json_or_yaml
from pptx_cli.core.manifest_store import template_copy_path
from pptx_cli.core.markdown import (
    ParsedParagraph,
    ParsedRun,
    looks_like_markdown,
    parse_markdown_paragraphs,
    parse_plain_text_paragraphs,
)
from pptx_cli.models.manifest import DeckSpec, LayoutContract, ManifestDocument, SlideSpec


class CompositionError(ValueError):
    code: str

    def __init__(self, code: str, message: str) -> None:
        super().__init__(message)
        self.code = code


_IMAGE_FIT_ALIASES = {
    "contain": "fit",
    "cover": "cover",
    "crop": "cover",
    "fill": "cover",
    "fit": "fit",
}

_MARKDOWN_HEADING_SPACE_AFTER_PT = {
    1: 12.0,
    2: 10.0,
    3: 8.0,
}
_MARKDOWN_HEADING_SPACE_BEFORE_PT = 6.0
_MARKDOWN_BODY_SPACE_AFTER_PT = 6.0
_MARKDOWN_LIST_SPACE_AFTER_PT = 2.0
_MARKDOWN_BLOCK_SPACE_BEFORE_PT = 6.0
_PPTX_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}
_TEXT_STYLE_TAGS = {
    "title": "titleStyle",
    "body": "bodyStyle",
    "other": "otherStyle",
}


def plan_output_change(output_path: Path, *, overwrite: bool) -> dict[str, str]:
    if output_path.exists() and not overwrite:
        raise CompositionError(
            "ERR_CONFLICT_OUTPUT_EXISTS",
            f"Output already exists: {output_path}. Use --overwrite to replace it.",
        )
    operation = "replace" if output_path.exists() else "create"
    return {
        "target": str(output_path),
        "operation": operation,
        "artifact_type": "pptx",
    }


def resolve_layout(manifest: ManifestDocument, layout_id: str) -> LayoutContract:
    for layout in manifest.layouts:
        if layout.id == layout_id or layout.name == layout_id or layout_id in layout.aliases:
            return layout
    raise CompositionError("ERR_VALIDATION_LAYOUT_UNKNOWN", f"Unknown layout: {layout_id}")


def parse_set_arguments(items: list[str]) -> dict[str, Any]:
    parsed: dict[str, Any] = {}
    for item in items:
        if "=" not in item:
            raise CompositionError("ERR_VALIDATION_SET_FORMAT", f"Invalid --set entry: {item}")
        key, raw_value = item.split("=", 1)
        parsed[key] = _load_inline_or_file_value(raw_value)
    return parsed


def _load_inline_or_file_value(raw_value: str) -> Any:
    if raw_value.startswith("@"):
        path = Path(raw_value[1:])
        if not path.exists():
            raise CompositionError("ERR_IO_NOT_FOUND", f"Referenced content file not found: {path}")
        if path.suffix.lower() in {".json", ".yaml", ".yml"}:
            return load_json_or_yaml(path)
        if path.suffix.lower() in {".md", ".txt"}:
            kind = "markdown-text" if path.suffix.lower() == ".md" else "text"
            return {
                "kind": kind,
                "value": path.read_text(encoding="utf-8"),
            }
        return {"kind": "image", "path": str(path)}

    if raw_value.startswith("{") or raw_value.startswith("["):
        try:
            return json.loads(raw_value)
        except json.JSONDecodeError:
            return raw_value
    return raw_value


def create_single_slide_spec(layout: str, content: dict[str, Any]) -> DeckSpec:
    return DeckSpec(slides=[SlideSpec(layout=layout, content=content)])


def build_presentation(manifest_dir: Path, manifest: ManifestDocument, spec: DeckSpec) -> Any:
    template_path = manifest_dir / manifest.template.stored_template_path
    if not template_path.exists():
        template_path = template_copy_path(manifest_dir)
    prs = Presentation(str(template_path))
    _remove_all_slides(prs)

    for slide_spec in spec.slides:
        layout_contract = resolve_layout(manifest, slide_spec.layout)
        slide_layout = prs.slide_layouts[layout_contract.source_layout_index]
        slide = prs.slides.add_slide(slide_layout)
        _populate_slide(slide, layout_contract, slide_spec.content)

    _apply_deck_metadata(prs, spec.metadata)
    return prs


def save_presentation(prs: Any, output_path: Path, *, overwrite: bool) -> None:
    if output_path.exists() and not overwrite:
        raise CompositionError(
            "ERR_CONFLICT_OUTPUT_EXISTS",
            f"Output already exists: {output_path}. Use --overwrite to replace it.",
        )
    output_path.parent.mkdir(parents=True, exist_ok=True)
    temp_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=output_path.suffix,
            dir=output_path.parent,
        ) as handle:
            temp_path = Path(handle.name)
        prs.save(str(temp_path))
        os.replace(temp_path, output_path)
    except PermissionError as exc:
        code = "ERR_CONFLICT_OUTPUT_EXISTS" if output_path.exists() else "ERR_IO_WRITE"
        if code == "ERR_CONFLICT_OUTPUT_EXISTS":
            message = f"Output file is in use or locked: {output_path}"
        else:
            message = f"Cannot write output file: {output_path}"
        raise CompositionError(code, message) from exc
    except OSError as exc:
        raise CompositionError("ERR_IO_WRITE", f"Cannot write output file: {output_path}") from exc
    finally:
        if temp_path is not None and temp_path.exists():
            temp_path.unlink(missing_ok=True)


def _remove_all_slides(prs: Any) -> None:
    slide_id_list = list(prs.slides._sldIdLst)
    for slide_id in slide_id_list:
        relationship_id = slide_id.rId
        prs.part.drop_rel(relationship_id)
        prs.slides._sldIdLst.remove(slide_id)


def _apply_deck_metadata(prs: Any, metadata: dict[str, Any]) -> None:
    core_properties = prs.core_properties
    title = metadata.get("title")
    author = metadata.get("author")
    if isinstance(title, str):
        core_properties.title = title
    if isinstance(author, str):
        core_properties.author = author


def _populate_slide(slide: Any, layout: LayoutContract, content: dict[str, Any]) -> None:
    expected = {placeholder.logical_name: placeholder for placeholder in layout.placeholders}
    unknown_keys = sorted(set(content) - set(expected))
    if unknown_keys:
        raise CompositionError(
            "ERR_VALIDATION_PLACEHOLDER_UNKNOWN",
            f"Unknown placeholders for layout {layout.id}: {', '.join(unknown_keys)}",
        )

    missing_required = [
        name
        for name, placeholder in expected.items()
        if placeholder.required and name not in content
    ]
    if missing_required:
        raise CompositionError(
            "ERR_VALIDATION_PLACEHOLDER_REQUIRED",
            f"Missing required placeholders for layout {layout.id}: {', '.join(missing_required)}",
        )

    filled_placeholder_idxs: set[int] = set()
    for key, value in content.items():
        placeholder = expected[key]
        shape = _find_slide_placeholder(slide, placeholder.placeholder_idx)
        if shape is None:
            raise CompositionError(
                "ERR_INTERNAL_PLACEHOLDER_MISSING",
                f"Placeholder {key} was not found on generated slide",
            )
        _apply_content_value(shape, placeholder.supported_content_types, value)
        filled_placeholder_idxs.add(placeholder.placeholder_idx)

    _cleanup_unused_placeholders(slide, layout, filled_placeholder_idxs)


def _find_slide_placeholder(slide: Any, placeholder_idx: int) -> Any | None:
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == placeholder_idx:
            return shape
    return None


def _apply_content_value(shape: Any, supported_types: list[str], value: Any) -> None:
    content = _normalize_content_value(value)
    kind = content["kind"]
    if kind not in supported_types:
        raise CompositionError(
            "ERR_VALIDATION_CONTENT_TYPE",
            f"Content type {kind!r} is not supported for this placeholder",
        )

    if kind in {"text", "markdown-text"}:
        _apply_text(shape, str(content["value"]), markdown=kind == "markdown-text")
        return
    if kind == "image":
        _apply_image(shape, content)
        return
    if kind == "table":
        _apply_table(shape, content)
        return
    if kind == "chart":
        _apply_chart(shape, content)
        return

    raise CompositionError("ERR_VALIDATION_CONTENT_TYPE", f"Unsupported content type: {kind}")


def _normalize_content_value(value: Any) -> dict[str, Any]:
    if isinstance(value, dict) and "kind" in value:
        return value
    if isinstance(value, str):
        if looks_like_markdown(value):
            return {"kind": "markdown-text", "value": value}
        return {"kind": "text", "value": value}
    if isinstance(value, (int, float, bool)):
        return {"kind": "text", "value": str(value)}
    if isinstance(value, list):
        return {"kind": "markdown-text", "value": "\n".join(f"- {item}" for item in value)}
    raise CompositionError("ERR_VALIDATION_CONTENT_TYPE", f"Unsupported content payload: {value!r}")


def _apply_text(shape: Any, text: str, *, markdown: bool) -> None:
    text_frame = shape.text_frame
    text_frame_state = _capture_text_frame_state(shape)
    text_frame.clear()
    _restore_text_frame_state(text_frame, text_frame_state)
    paragraphs = parse_markdown_paragraphs(text) if markdown else parse_plain_text_paragraphs(text)
    list_level_offset = _first_markdown_list_level(shape) if markdown else 0

    previous: ParsedParagraph | None = None
    for index, parsed in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        _apply_markdown_list_level(paragraph, parsed, list_level_offset)
        _write_paragraph_runs(paragraph, parsed)
        if markdown:
            _apply_markdown_paragraph_format(paragraph, parsed, previous, is_first=index == 0)
        previous = parsed


def _capture_text_frame_state(shape: Any) -> dict[str, Any]:
    text_frame = shape.text_frame
    return {
        "vertical_anchor": _resolve_text_frame_value(shape, "vertical_anchor"),
        "word_wrap": _resolve_text_frame_value(shape, "word_wrap"),
        "auto_size": _resolve_text_frame_value(shape, "auto_size"),
        "margin_top": text_frame.margin_top,
        "margin_bottom": text_frame.margin_bottom,
        "margin_left": text_frame.margin_left,
        "margin_right": text_frame.margin_right,
    }


def _restore_text_frame_state(text_frame: Any, state: dict[str, Any]) -> None:
    text_frame.vertical_anchor = state["vertical_anchor"]
    text_frame.word_wrap = state["word_wrap"]
    text_frame.auto_size = state["auto_size"]
    text_frame.margin_top = state["margin_top"]
    text_frame.margin_bottom = state["margin_bottom"]
    text_frame.margin_left = state["margin_left"]
    text_frame.margin_right = state["margin_right"]


def _resolve_text_frame_value(shape: Any, attribute_name: str) -> Any:
    value = getattr(shape.text_frame, attribute_name)
    if value is not None:
        return value

    layout_placeholder = _find_layout_placeholder(shape)
    if layout_placeholder is None or not hasattr(layout_placeholder, "text_frame"):
        return value
    return getattr(layout_placeholder.text_frame, attribute_name)


def _find_layout_placeholder(shape: Any) -> Any | None:
    if not getattr(shape, "is_placeholder", False):
        return None

    placeholder_idx = int(shape.placeholder_format.idx)
    for layout_shape in shape.part.slide_layout.placeholders:
        if int(layout_shape.placeholder_format.idx) == placeholder_idx:
            return layout_shape
    return None


def _apply_markdown_list_level(
    paragraph: Any,
    parsed: ParsedParagraph,
    list_level_offset: int,
) -> None:
    if parsed.kind not in {"bullet", "ordered"} or parsed.level is None:
        return

    paragraph.level = min(list_level_offset + parsed.level, 8)
    _set_paragraph_bullet_mode(paragraph, enabled=parsed.kind == "bullet")


def _set_paragraph_bullet_mode(paragraph: Any, *, enabled: bool) -> None:
    paragraph_properties = _paragraph_properties(paragraph)
    for tag_name in ("a:buNone", "a:buChar", "a:buAutoNum"):
        child = paragraph_properties.find(qn(tag_name))
        if child is not None:
            paragraph_properties.remove(child)
    if not enabled:
        paragraph_properties.insert(0, OxmlElement("a:buNone"))


def _paragraph_properties(paragraph: Any) -> Any:
    paragraph_properties = paragraph._element.find(qn("a:pPr"))
    if paragraph_properties is not None:
        return paragraph_properties

    paragraph_properties = OxmlElement("a:pPr")
    paragraph._element.insert(0, paragraph_properties)
    return paragraph_properties


def _first_markdown_list_level(shape: Any) -> int:
    style_root = _text_style_root(shape)
    if style_root is None:
        return 0

    for level in range(0, 5):
        if _paragraph_style_has_bullet(_paragraph_style_for_level(style_root, level)):
            return level
    return 0


def _text_style_root(shape: Any) -> Any | None:
    tag_name = _TEXT_STYLE_TAGS[_placeholder_text_style_bucket(shape)]
    return shape.part.slide_layout.slide_master.element.find(
        f".//p:{tag_name}",
        namespaces=_PPTX_NS,
    )


def _placeholder_text_style_bucket(shape: Any) -> str:
    if not getattr(shape, "is_placeholder", False):
        return "body"

    placeholder_type = shape.placeholder_format.type.name.lower()
    if placeholder_type in {"title", "center_title", "vertical_title"}:
        return "title"
    if placeholder_type in {"body", "object", "content", "text", "subtitle"}:
        return "body"
    return "other"


def _paragraph_style_for_level(style_root: Any, level: int) -> Any | None:
    if level == 0:
        paragraph_style = style_root.find("./a:lvl1pPr", namespaces=_PPTX_NS)
        if paragraph_style is not None:
            return paragraph_style
        return style_root.find("./a:defPPr", namespaces=_PPTX_NS)

    return style_root.find(f"./a:lvl{level + 1}pPr", namespaces=_PPTX_NS)


def _paragraph_style_has_bullet(paragraph_style: Any | None) -> bool:
    if paragraph_style is None:
        return False
    if paragraph_style.find(qn("a:buNone")) is not None:
        return False
    return (
        paragraph_style.find(qn("a:buChar")) is not None
        or paragraph_style.find(qn("a:buAutoNum")) is not None
    )


def _write_paragraph_runs(paragraph: Any, parsed: ParsedParagraph) -> None:
    first_run_spec = parsed.runs[0]
    paragraph.text = first_run_spec.text
    _apply_run_format(paragraph.runs[0], first_run_spec)

    for run_spec in parsed.runs[1:]:
        run = paragraph.add_run()
        run.text = run_spec.text
        _apply_run_format(run, run_spec)


def _apply_run_format(run: Any, parsed: ParsedRun) -> None:
    if parsed.bold:
        run.font.bold = True
    if parsed.italic:
        run.font.italic = True
    if parsed.code:
        run.font.name = "Courier New"


def _apply_markdown_paragraph_format(
    paragraph: Any,
    parsed: ParsedParagraph,
    previous: ParsedParagraph | None,
    *,
    is_first: bool,
) -> None:
    if parsed.kind == "heading":
        _apply_heading_runs(paragraph)
        if not is_first:
            paragraph.space_before = Pt(_MARKDOWN_HEADING_SPACE_BEFORE_PT)
        paragraph.space_after = Pt(
            _MARKDOWN_HEADING_SPACE_AFTER_PT.get(parsed.heading_level or 3, 8.0)
        )
        return

    if _starts_new_markdown_block(previous, parsed):
        paragraph.space_before = Pt(_MARKDOWN_BLOCK_SPACE_BEFORE_PT)

    if parsed.kind in {"bullet", "ordered"}:
        paragraph.space_after = Pt(_MARKDOWN_LIST_SPACE_AFTER_PT)
        return

    paragraph.space_after = Pt(_MARKDOWN_BODY_SPACE_AFTER_PT)


def _apply_heading_runs(paragraph: Any) -> None:
    for run in paragraph.runs:
        run.font.bold = True


def _starts_new_markdown_block(
    previous: ParsedParagraph | None,
    current: ParsedParagraph,
) -> bool:
    if previous is None:
        return False
    if current.kind in {"bullet", "ordered"}:
        return previous.kind != current.kind
    return previous.kind in {"bullet", "ordered"}


def _cleanup_unused_placeholders(
    slide: Any,
    layout: LayoutContract,
    filled_placeholder_idxs: set[int],
) -> None:
    expected_by_idx = {
        placeholder.placeholder_idx: placeholder for placeholder in layout.placeholders
    }
    for shape in list(slide.placeholders):
        placeholder_idx = int(shape.placeholder_format.idx)
        placeholder = expected_by_idx.get(placeholder_idx)
        if placeholder is None:
            continue
        if placeholder_idx in filled_placeholder_idxs:
            continue
        if placeholder.required:
            continue
        shape._element.getparent().remove(shape._element)


def _apply_image(shape: Any, content: dict[str, Any]) -> None:
    image_path = Path(str(content["path"]))
    if not image_path.exists():
        raise CompositionError("ERR_IO_NOT_FOUND", f"Image not found: {image_path}")

    placeholder_size = (int(shape.width), int(shape.height))
    image_fit = _normalize_image_fit(content.get("image_fit"))
    picture, image_size = _insert_placeholder_picture(shape, image_path)
    _apply_image_crop(picture, image_size, placeholder_size, image_fit)


def _apply_table(shape: Any, content: dict[str, Any]) -> None:
    columns = content.get("columns", [])
    rows = content.get("rows", [])
    if not isinstance(columns, list) or not isinstance(rows, list):
        raise CompositionError(
            "ERR_VALIDATION_TABLE_PAYLOAD",
            "Table payload requires 'columns' and 'rows' lists",
        )

    row_count = len(rows) + (1 if columns else 0)
    col_count = len(columns) if columns else (len(rows[0]) if rows else 0)
    if row_count <= 0 or col_count <= 0:
        raise CompositionError(
            "ERR_VALIDATION_TABLE_PAYLOAD",
            "Table payload must contain at least one row and one column",
        )

    graphic_frame = _insert_placeholder_table(shape, row_count, col_count)
    table = graphic_frame.table
    row_offset = 0
    if columns:
        for column_index, heading in enumerate(columns):
            table.cell(0, column_index).text = str(heading)
        row_offset = 1
    for row_index, row in enumerate(rows, start=row_offset):
        for column_index, cell in enumerate(row):
            table.cell(row_index, column_index).text = str(cell)


def _apply_chart(shape: Any, content: dict[str, Any]) -> None:
    chart_data = ChartData()
    categories = content.get("categories", [])
    series = content.get("series", [])
    chart_type_name = str(content.get("chart_type", "column_clustered")).upper()
    if not isinstance(categories, list) or not isinstance(series, list):
        raise CompositionError(
            "ERR_VALIDATION_CHART_PAYLOAD",
            "Chart payload requires 'categories' and 'series' lists",
        )

    chart_data.categories = categories
    for series_entry in series:
        if not isinstance(series_entry, dict):
            raise CompositionError(
                "ERR_VALIDATION_CHART_PAYLOAD",
                "Each chart series must be an object",
            )
        chart_data.add_series(
            str(series_entry.get("name", "Series")),
            series_entry.get("values", []),
        )

    if not hasattr(XL_CHART_TYPE, chart_type_name):
        raise CompositionError(
            "ERR_VALIDATION_CHART_PAYLOAD",
            f"Unsupported chart_type: {chart_type_name.lower()}",
        )
    chart_type = getattr(XL_CHART_TYPE, chart_type_name)
    _insert_placeholder_chart(shape, chart_type, chart_data)


def _normalize_image_fit(value: Any) -> str:
    if value is None:
        return "fit"
    if not isinstance(value, str):
        raise CompositionError(
            "ERR_VALIDATION_IMAGE_FIT",
            "image_fit must be one of: fit, contain, cover, fill, crop",
        )
    normalized = value.strip().lower()
    if normalized not in _IMAGE_FIT_ALIASES:
        raise CompositionError(
            "ERR_VALIDATION_IMAGE_FIT",
            "image_fit must be one of: fit, contain, cover, fill, crop",
        )
    return _IMAGE_FIT_ALIASES[normalized]


def _insert_placeholder_picture(
    shape: Any,
    image_path: Path,
) -> tuple[PlaceholderPicture, tuple[int, int]]:
    image_part, relationship_id = shape.part.get_or_add_image_part(str(image_path))
    image_size = image_part._px_size
    picture_element = CT_Picture.new_ph_pic(
        shape.shape_id,
        shape.name,
        image_part.desc,
        relationship_id,
    )
    _replace_placeholder_with_preserved_order(shape, picture_element)
    return PlaceholderPicture(picture_element, shape._parent), image_size


def _insert_placeholder_table(shape: Any, rows: int, cols: int) -> PlaceholderGraphicFrame:
    graphic_frame = CT_GraphicalObjectFrame.new_table_graphicFrame(
        shape.shape_id,
        shape.name,
        rows,
        cols,
        shape.left,
        shape.top,
        shape.width,
        Emu(rows * 370840),
    )
    _replace_placeholder_with_preserved_order(shape, graphic_frame)
    return PlaceholderGraphicFrame(graphic_frame, shape._parent)


def _insert_placeholder_chart(
    shape: Any,
    chart_type: XL_CHART_TYPE,
    chart_data: ChartData,
) -> PlaceholderGraphicFrame:
    relationship_id = shape.part.add_chart_part(chart_type, chart_data)
    graphic_frame = CT_GraphicalObjectFrame.new_chart_graphicFrame(
        shape.shape_id,
        shape.name,
        relationship_id,
        shape.left,
        shape.top,
        shape.width,
        shape.height,
    )
    _replace_placeholder_with_preserved_order(shape, graphic_frame)
    return PlaceholderGraphicFrame(graphic_frame, shape._parent)


def _replace_placeholder_with_preserved_order(shape: Any, replacement: Any) -> None:
    shape_tree = shape._element.getparent()
    original_index = list(shape_tree).index(shape._element)
    shape._replace_placeholder_with(replacement)
    current_index = list(shape_tree).index(replacement)
    if current_index == original_index:
        return
    shape_tree.remove(replacement)
    shape_tree.insert(original_index, replacement)


def _apply_image_crop(
    picture: PlaceholderPicture,
    image_size: tuple[int, int],
    placeholder_size: tuple[int, int],
    image_fit: str,
) -> None:
    if image_fit == "cover":
        picture._pic.crop_to_fit(image_size, placeholder_size)
        return

    image_width, image_height = image_size
    placeholder_width, placeholder_height = placeholder_size
    if image_width <= 0 or image_height <= 0:
        return
    if placeholder_width <= 0 or placeholder_height <= 0:
        return

    picture.crop_left = 0.0
    picture.crop_top = 0.0
    picture.crop_right = 0.0
    picture.crop_bottom = 0.0

    view_ratio = placeholder_width / placeholder_height
    image_ratio = image_width / image_height

    if image_ratio > view_ratio:
        padding = (image_ratio / view_ratio - 1.0) / 2.0
        picture.crop_top = -padding
        picture.crop_bottom = -padding
        return
    if image_ratio < view_ratio:
        padding = (view_ratio / image_ratio - 1.0) / 2.0
        picture.crop_left = -padding
        picture.crop_right = -padding
